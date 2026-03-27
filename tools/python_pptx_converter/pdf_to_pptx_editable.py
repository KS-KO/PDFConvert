import sys, os, argparse, asyncio, io, fitz
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw

# Optional: winsdk for Windows OCR
HAS_WINSDK = False
try:
    import winsdk.windows.graphics.imaging as imaging
    import winsdk.windows.media.ocr as ocr
    import winsdk.windows.storage.streams as streams
    HAS_WINSDK = True
except ImportError:
    pass

FONT_MAP = {"Arial":"Arial","TimesNewRoman":"Times New Roman","MalgunGothic":"Malgun Gothic","NanumGothic":"Nanum Gothic"}

async def perform_windows_ocr(image_bytes):
    if not HAS_WINSDK: return []
    try:
        stream = streams.InMemoryRandomAccessStream()
        writer = streams.DataWriter(stream)
        writer.write_bytes(image_bytes)
        await writer.store_async(); await writer.flush_async(); stream.seek(0)
        decoder = await imaging.BitmapDecoder.create_async(stream)
        bitmap = await decoder.get_software_bitmap_async()
        
        # Prefer Korean, fall back to anything
        engine = ocr.OcrEngine.try_create_from_user_profile_languages()
        if not engine:
             # Try first available
             langs = ocr.OcrEngine.available_recognizer_languages
             if langs: engine = ocr.OcrEngine.try_create_from_language(langs[0])
        
        if not engine: return []
        
        result = await engine.recognize_async(bitmap)
        return [{"text":w.text, "bbox":[w.bounding_rect.x, w.bounding_rect.y, w.bounding_rect.x+w.bounding_rect.width, w.bounding_rect.y+w.bounding_rect.height]} for l in result.lines for w in l.words]
    except Exception as e:
        print(f"OCR Error: {e}")
        return []

def get_colors(image, bbox):
    try:
        w, h = image.size
        # Unpack and Ensure Integers
        x0, y0, x1, y1 = [int(round(v)) for v in bbox]
        
        # Guard - avoid crash if bbox is zero-width or out of bounds
        x0 = max(0, min(w-1, x0))
        y0 = max(0, min(h-1, y0))
        x1 = max(x0+1, min(w, x1))
        y1 = max(y0+1, min(h, y1))

        if x1 <= x0 or y1 <= y0: return (255,255,255), (0,0,0)

        # Sample background (ring)
        pixels = []
        p = 4
        # Range around the box
        for x in range(max(0, x0-p), min(w, x1+p), 2):
            # Top edge
            for y in range(max(0, y0-p), max(0, y0), 2): pixels.append(image.getpixel((x,y)))
            # Bottom edge
            for y in range(min(h, y1), min(h, y1+p), 2): pixels.append(image.getpixel((x,y)))
        
        bg = (sum(p[0] for p in pixels)//len(pixels), sum(p[1] for p in pixels)//len(pixels), sum(p[2] for p in pixels)//len(pixels)) if pixels else (255,255,255)
        
        # Sample font
        f_pixels = []
        for x in range(x0, x1):
            for y in range(y0, y1):
                pix = image.getpixel((x,y))
                # Distance from BG
                if sum((pix[i]-bg[i])**2 for i in range(3))**0.5 > 35: f_pixels.append(pix)
        
        if not f_pixels: return bg, (0,0,0)
        
        # Heuristic: sort by luminance
        f_pixels.sort(key=lambda x: sum(x), reverse=(sum(bg)/3.0 > 128))
        sample = f_pixels[:len(f_pixels)//4 or 1]
        font = (sum(p[0] for p in sample)//len(sample), sum(p[1] for p in sample)//len(sample), sum(p[2] for p in sample)//len(sample))
        return bg, font
    except Exception as e:
        print(f"Color Sample Error: {e}")
        return (255,255,255), (0,0,0)

async def pdf_to_pptx_hybrid(pdf_path, output_dir):
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"Failed to open PDF: {e}")
        return

    prs = Presentation()
    for page in doc:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            if page.number == 0: 
                prs.slide_width, prs.slide_height = Inches(page.rect.width/72), Inches(page.rect.height/72)
            
            # 1. Native text
            items = []
            try:
                 text_dict = page.get_text("dict")
                 for b in text_dict.get("blocks", []):
                    if b["type"] == 0:
                        for l in b["lines"]:
                            for s in l["spans"]:
                                items.append({"text":s["text"], "bbox":s["bbox"], "size":s["size"], "color":s["color"], "font":s["font"], "native":True})
            except: pass

            # 2. OCR for missing parts (High Res)
            zoom = 3.0
            try:
                pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                img_data = pix.tobytes("png")
                ocr_res = await perform_windows_ocr(img_data)
                pil_img = Image.open(io.BytesIO(img_data)).convert("RGB")
                draw = ImageDraw.Draw(pil_img)
                
                for o in ocr_res:
                    ob = [v/zoom for v in o["bbox"]]
                    # Filter duplicates (Overlap > 40%)
                    is_duplicate = False
                    for n in items:
                        if not n.get("native"): continue
                        nb = n["bbox"]
                        inter_x = max(0, min(ob[2], nb[2]) - max(ob[0], nb[0]))
                        inter_y = max(0, min(ob[3], nb[3]) - max(ob[1], nb[1]))
                        if inter_x * inter_y > (ob[2]-ob[0])*(ob[3]-ob[1]) * 0.4:
                            is_duplicate = True
                            break
                    
                    if not is_duplicate:
                        bg, font = get_colors(pil_img, o["bbox"])
                        items.append({"text":o["text"], "bbox":ob, "size":(o["bbox"][3]-o["bbox"][1])/zoom, "rgb":font, "native":False})
                        draw.rectangle([v-1 for v in o["bbox"]], fill=bg) # Erase from img
                
                # Save Background
                bg_io = io.BytesIO(); pil_img.save(bg_io, format="PNG"); bg_io.seek(0)
                slide.shapes.add_picture(bg_io, 0, 0, width=prs.slide_width, height=prs.slide_height)
            except Exception as e:
                print(f"Background/OCR Error on page {page.number}: {e}")

            # 3. Add Text Boxes
            for it in items:
                try:
                    tx = it.get("text", "").strip()
                    if not tx: continue
                    b = it["bbox"]
                    # Guard against invalid dimensions
                    w = max(0.1, (b[2]-b[0])/72)
                    h = max(0.1, (b[3]-b[1])/72)
                    shape = slide.shapes.add_textbox(Inches(b[0]/72), Inches(b[1]/72), Inches(w), Inches(h))
                    
                    # Try to add paragraph and run safely
                    if not shape.text_frame.paragraphs:
                        p = shape.text_frame.add_paragraph()
                    else:
                        p = shape.text_frame.paragraphs[0]
                    
                    run = p.add_run()
                    run.text = tx
                    
                    # Font Size
                    size = it.get("size", 10)
                    if size <= 0: size = 10
                    run.font.size = Pt(size)
                    
                    # Font Color
                    if "rgb" in it and it["rgb"]:
                        run.font.color.rgb = RGBColor(*[int(v) for v in it["rgb"]])
                    elif "color" in it and it["color"] is not None: 
                        c = int(it["color"])
                        run.font.color.rgb = RGBColor((c>>16)&0xFF, (c>>8)&0xFF, c&0xFF)
                except Exception as ex:
                    print(f"Textbox skip on page {page.number}: {ex}")
        except Exception as e:
            print(f"Slide creation failed on page {page.number}: {e}")
            
    # Correct Suffix: _editable.pptx (to match C# OutputWriter)
    out_p = os.path.join(output_dir, os.path.splitext(os.path.basename(pdf_path))[0] + "_editable.pptx")
    print(f"Saving to: {out_p}")
    try:
        prs.save(out_p)
        print(f"DONE: {out_p}")
    except Exception as e:
        print(f"Save Failed: {e}")
        sys.exit(1)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(); parser.add_argument("--pdf", required=True); parser.add_argument("--outdir", required=True)
    args = parser.parse_args()
    print(f"Script started with PDF: {args.pdf}, OutDir: {args.outdir}")
    asyncio.run(pdf_to_pptx_hybrid(args.pdf, args.outdir))
